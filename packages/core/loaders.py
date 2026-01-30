"""Document loaders for various file formats."""
import hashlib
from pathlib import Path
from typing import List, Tuple, Optional
import json

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd
from bs4 import BeautifulSoup
import markdown

from packages.core.logging_config import setup_logging

logger = setup_logging(__name__)


class Section:
    """Document section with stable source_ref."""
    
    def __init__(self, source_ref: str, content: str, metadata: Optional[dict] = None):
        self.source_ref = source_ref
        self.content = content
        self.metadata = metadata or {}
    
    def to_dict(self) -> dict:
        return {
            "source_ref": self.source_ref,
            "content": self.content,
            "metadata": self.metadata
        }


def compute_file_sha256(filepath: str) -> str:
    """Compute SHA256 hash of a file."""
    sha256_hash = hashlib.sha256()
    with open(filepath, "rb") as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()


def load_pdf(filepath: str) -> List[Section]:
    """
    Load PDF and extract sections by page.
    
    Args:
        filepath: Path to PDF file
        
    Returns:
        List of Section objects with source_ref='page=N'
    """
    sections = []
    
    try:
        doc = fitz.open(filepath)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            if text.strip():
                source_ref = f"page={page_num + 1}"
                sections.append(Section(source_ref, text.strip()))
        
        doc.close()
        logger.info(f"Loaded PDF: {filepath}, {len(sections)} pages")
    except Exception as e:
        logger.error(f"Error loading PDF {filepath}: {e}")
        raise
    
    return sections


def load_docx(filepath: str) -> List[Section]:
    """
    Load DOCX and extract sections by heading or paragraph groups.
    
    Args:
        filepath: Path to DOCX file
        
    Returns:
        List of Section objects
    """
    sections = []
    
    try:
        doc = DocxDocument(filepath)
        current_heading = "document_start"
        current_content = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Check if paragraph is a heading
            if para.style.name.startswith('Heading'):
                # Save previous section
                if current_content:
                    content_text = '\n'.join(current_content)
                    sections.append(Section(f"heading={current_heading}", content_text))
                
                # Start new section
                current_heading = text[:100]  # Truncate long headings
                current_content = []
            else:
                current_content.append(text)
        
        # Add final section
        if current_content:
            content_text = '\n'.join(current_content)
            sections.append(Section(f"heading={current_heading}", content_text))
        
        # If no sections created, create one for the whole document
        if not sections:
            all_text = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            if all_text.strip():
                sections.append(Section("section=all", all_text.strip()))
        
        logger.info(f"Loaded DOCX: {filepath}, {len(sections)} sections")
    except Exception as e:
        logger.error(f"Error loading DOCX {filepath}: {e}")
        raise
    
    return sections


def load_pptx(filepath: str) -> List[Section]:
    """
    Load PPTX and extract sections by slide.
    
    Args:
        filepath: Path to PPTX file
        
    Returns:
        List of Section objects with source_ref='slide=N'
    """
    sections = []
    
    try:
        prs = Presentation(filepath)
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            
            # Extract title
            if slide.shapes.title:
                texts.append(slide.shapes.title.text)
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            
            content = '\n'.join(texts).strip()
            if content:
                source_ref = f"slide={slide_num}"
                sections.append(Section(source_ref, content))
        
        logger.info(f"Loaded PPTX: {filepath}, {len(sections)} slides")
    except Exception as e:
        logger.error(f"Error loading PPTX {filepath}: {e}")
        raise
    
    return sections


def load_xlsx(filepath: str) -> List[Section]:
    """
    Load XLSX and extract sections by sheet.
    
    Args:
        filepath: Path to XLSX file
        
    Returns:
        List of Section objects with source_ref='sheet=NAME'
    """
    sections = []
    
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Convert sheet to text representation
            rows = []
            for row in ws.iter_rows(values_only=True):
                # Filter out empty rows
                row_values = [str(cell) if cell is not None else '' for cell in row]
                if any(val.strip() for val in row_values):
                    rows.append('\t'.join(row_values))
            
            content = '\n'.join(rows).strip()
            if content:
                source_ref = f"sheet={sheet_name}"
                sections.append(Section(source_ref, content))
        
        wb.close()
        logger.info(f"Loaded XLSX: {filepath}, {len(sections)} sheets")
    except Exception as e:
        logger.error(f"Error loading XLSX {filepath}: {e}")
        raise
    
    return sections


def load_html(filepath: str) -> List[Section]:
    """
    Load HTML and extract sections by heading.
    
    Args:
        filepath: Path to HTML file
        
    Returns:
        List of Section objects
    """
    sections = []
    
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(['script', 'style']):
            script.decompose()
        
        # Try to extract by headings
        current_heading = "document_start"
        current_content = []
        
        for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'div', 'li']):
            text = element.get_text().strip()
            if not text:
                continue
            
            if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                # Save previous section
                if current_content:
                    content_text = '\n'.join(current_content)
                    sections.append(Section(f"heading={current_heading}", content_text))
                
                # Start new section
                current_heading = text[:100]
                current_content = []
            else:
                current_content.append(text)
        
        # Add final section
        if current_content:
            content_text = '\n'.join(current_content)
            sections.append(Section(f"heading={current_heading}", content_text))
        
        # If no sections, create one for all text
        if not sections:
            text = soup.get_text()
            clean_text = '\n'.join([line.strip() for line in text.split('\n') if line.strip()])
            if clean_text:
                sections.append(Section("section=all", clean_text))
        
        logger.info(f"Loaded HTML: {filepath}, {len(sections)} sections")
    except Exception as e:
        logger.error(f"Error loading HTML {filepath}: {e}")
        raise
    
    return sections


def load_markdown(filepath: str) -> List[Section]:
    """
    Load Markdown and extract sections by heading.
    
    Args:
        filepath: Path to Markdown file
        
    Returns:
        List of Section objects
    """
    sections = []
    
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Convert markdown to HTML first, then extract sections
        html = markdown.markdown(content)
        soup = BeautifulSoup(html, 'html.parser')
        
        current_heading = "document_start"
        current_content = []
        
        for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li']):
            text = element.get_text().strip()
            if not text:
                continue
            
            if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                # Save previous section
                if current_content:
                    content_text = '\n'.join(current_content)
                    sections.append(Section(f"heading={current_heading}", content_text))
                
                # Start new section
                current_heading = text[:100]
                current_content = []
            else:
                current_content.append(text)
        
        # Add final section
        if current_content:
            content_text = '\n'.join(current_content)
            sections.append(Section(f"heading={current_heading}", content_text))
        
        # If no sections, create one for all text
        if not sections:
            sections.append(Section("section=all", content.strip()))
        
        logger.info(f"Loaded Markdown: {filepath}, {len(sections)} sections")
    except Exception as e:
        logger.error(f"Error loading Markdown {filepath}: {e}")
        raise
    
    return sections


def load_txt(filepath: str) -> List[Section]:
    """
    Load plain text file.
    
    Args:
        filepath: Path to text file
        
    Returns:
        List of Section objects
    """
    sections = []
    
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read().strip()
        
        if content:
            filename = Path(filepath).stem
            sections.append(Section(f"section={filename}", content))
        
        logger.info(f"Loaded TXT: {filepath}")
    except Exception as e:
        logger.error(f"Error loading TXT {filepath}: {e}")
        raise
    
    return sections


def load_document(filepath: str) -> List[Section]:
    """
    Load document based on file extension.
    
    Args:
        filepath: Path to document
        
    Returns:
        List of Section objects
    """
    ext = Path(filepath).suffix.lower()
    
    loaders = {
        '.pdf': load_pdf,
        '.docx': load_docx,
        '.pptx': load_pptx,
        '.xlsx': load_xlsx,
        '.xls': load_xlsx,
        '.html': load_html,
        '.htm': load_html,
        '.md': load_markdown,
        '.txt': load_txt,
    }
    
    loader = loaders.get(ext)
    if not loader:
        raise ValueError(f"Unsupported file type: {ext}")
    
    return loader(filepath)
